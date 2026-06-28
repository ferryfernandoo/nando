#!/usr/bin/env python3
"""
Secure PowerPoint Generator
- Input validation & sanitization
- Resource limits (timeout, file size, slide count)
- Error handling
- Safe execution environment
"""

import re
import sys
import json
import os
import tempfile
import threading
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ====== PPT TEMPLATE STYLES =====
TEMPLATE_STYLES = {
    'classic': {
        'background_color': RGBColor(255, 255, 255),
        'accent_color': RGBColor(0, 102, 204),
        'title_color': RGBColor(0, 51, 102),
        'content_color': RGBColor(64, 64, 64),
        'panel_color': RGBColor(255, 255, 255),
        'accent_bar': True,
    },
    'modern': {
        'background_color': RGBColor(15, 23, 42),
        'accent_color': RGBColor(59, 130, 246),
        'title_color': RGBColor(255, 255, 255),
        'content_color': RGBColor(226, 232, 240),
        'panel_color': RGBColor(30, 41, 59),
        'accent_bar': True,
    },
    'bold': {
        'background_color': RGBColor(255, 112, 67),
        'accent_color': RGBColor(255, 255, 255),
        'title_color': RGBColor(255, 255, 255),
        'content_color': RGBColor(255, 255, 255),
        'panel_color': RGBColor(255, 178, 128),
        'accent_bar': False,
    },
    'minimal': {
        'background_color': RGBColor(249, 250, 251),
        'accent_color': RGBColor(100, 116, 139),
        'title_color': RGBColor(15, 23, 42),
        'content_color': RGBColor(51, 65, 85),
        'panel_color': RGBColor(255, 255, 255),
        'accent_bar': True,
    },
}

# ============ SECURITY CONFIG ============
TIMEOUT_SECONDS = 30
MAX_FILE_SIZE_MB = 50
MAX_SLIDES = 100
MAX_TEXT_LENGTH = 10000
MAX_TITLE_LENGTH = 200
ALLOWED_FORMATS = ['pptx']
OUTPUT_DIR = Path(__file__).parent / 'temp_ppt'

# ============ SECURITY LIMITS ============
def validate_input(data):
    """Validate and sanitize input data"""
    errors = []
    
    # Check required fields
    if not isinstance(data, dict):
        return False, ["Invalid input format"]
    
    # Validate title
    title = data.get('title', '').strip()
    if not title:
        errors.append("Title required")
    elif len(title) > MAX_TITLE_LENGTH:
        errors.append(f"Title too long (max {MAX_TITLE_LENGTH} chars)")
    
    # Validate slides
    slides = data.get('slides', [])
    if not isinstance(slides, list):
        errors.append("Slides must be an array")
    elif len(slides) > MAX_SLIDES:
        errors.append(f"Too many slides (max {MAX_SLIDES})")
    elif len(slides) == 0:
        errors.append("At least 1 slide required")
    
    # Validate each slide
    for i, slide in enumerate(slides):
        if not isinstance(slide, dict):
            errors.append(f"Slide {i+1}: Invalid format")
            continue
        
        slide_title = slide.get('title', '').strip()
        if len(slide_title) > MAX_TITLE_LENGTH:
            errors.append(f"Slide {i+1}: Title too long")
        
        content = slide.get('content', '').strip()
        if len(content) > MAX_TEXT_LENGTH:
            errors.append(f"Slide {i+1}: Content too long (max {MAX_TEXT_LENGTH} chars)")
    
    return len(errors) == 0, errors

def sanitize_text(text, max_length=MAX_TEXT_LENGTH):
    """Remove potentially harmful characters and limit length"""
    if not isinstance(text, str):
        text = str(text)
    
    # Remove null bytes and control characters
    text = ''.join(c for c in text if c.isprintable() or c in '\n\t\r')
    
    # Limit length
    return text[:max_length].strip()


def parse_markdown_table(text):
    """Detect markdown-style table blocks and return rows if present."""
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if len(lines) < 2:
        return None

    rows = []
    for line in lines:
        if '|' not in line:
            return None
        clean = line
        if clean.startswith('|'):
            clean = clean[1:]
        if clean.endswith('|'):
            clean = clean[:-1]
        row_cells = [cell.strip() for cell in clean.split('|')]
        rows.append(row_cells)

    if len(rows) > 1 and all(re.fullmatch(r':?-{2,}:?', cell) for cell in rows[1]):
        rows.pop(1)

    if len(rows) < 2:
        return None

    max_cols = max(len(r) for r in rows)
    normalized = [r + [''] * (max_cols - len(r)) for r in rows]
    return normalized


def add_text_box(slide, x, y, cx, cy, text, font_size, font_color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(cx), Inches(cy))
    text_frame = box.text_frame
    text_frame.word_wrap = True
    text_frame.clear()
    text_frame.text = sanitize_text(text)
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = font_color
        paragraph.alignment = align
    return box


def add_markdown_table(slide, rows, left, top, width, height, style):
    if not rows:
        return
    cols = len(rows[0])
    table = slide.shapes.add_table(len(rows), cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    for r_index, row in enumerate(rows):
        for c_index, cell_text in enumerate(row):
            cell = table.cell(r_index, c_index)
            cell.text = sanitize_text(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(16 if r_index else 18)
                    run.font.bold = (r_index == 0)
                    run.font.color.rgb = style['title_color'] if r_index == 0 else style['content_color']
            if r_index == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = style['accent_color']
    return table


def render_slide_content(slide, content, style):
    table_rows = parse_markdown_table(content)
    if table_rows:
        add_markdown_table(slide, table_rows, 0.6, 1.6, 8.8, 4.0, style)
        return

    add_text_box(slide, 0.6, 1.3, 8.8, 5.4, content or '', 18, style['content_color'], bold=False, align=PP_ALIGN.LEFT)

def timeout_handler():
    """Handle timeout - raises exception in main thread"""
    # This is called by timer thread, but we handle timeout in main with a flag
    pass

def setup_output_dir():
    """Setup temporary output directory"""
    try:
        OUTPUT_DIR.mkdir(exist_ok=True)
        # Clean old files (>1 hour)
        for f in OUTPUT_DIR.glob('*.pptx'):
            if (datetime.now() - datetime.fromtimestamp(f.stat().st_mtime)).seconds > 3600:
                f.unlink()
    except Exception as e:
        raise Exception(f"Failed to setup output directory: {e}")


def get_template_style(name):
    return TEMPLATE_STYLES.get(name if isinstance(name, str) else 'classic', TEMPLATE_STYLES['classic'])


def add_slide_with_content(prs, slide_title, content, style):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if style['background_color']:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = style['background_color']
        bg.line.fill.background()

    if style['accent_bar']:
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.45))
        accent.fill.solid()
        accent.fill.fore_color.rgb = style['accent_color']
        accent.line.fill.background()

    if style['panel_color'] and style['panel_color'] != style['background_color']:
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.35), Inches(8.8), Inches(4.85))
        panel.fill.solid()
        panel.fill.fore_color.rgb = style['panel_color']
        panel.line.fill.background()

    title_y = 0.7 if style['accent_bar'] else 0.5
    add_text_box(slide, 0.6, title_y, 8.8, 0.9, slide_title or 'Slide', 32, style['title_color'], bold=True, align=PP_ALIGN.LEFT)
    render_slide_content(slide, content or '', style)


def apply_title_slide_style(slide, title, subtitle, style):
    if style['background_color']:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide.part.slide_width, slide.part.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = style['background_color']
        bg.line.fill.background()

    if style['accent_bar']:
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide.part.slide_width, Inches(0.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = style['accent_color']
        accent.line.fill.background()

    add_text_box(slide, 0.6, 1.0, 8.8, 1.2, title or 'Untitled', 48, style['title_color'], bold=True, align=PP_ALIGN.LEFT)
    add_text_box(slide, 0.6, 2.4, 8.8, 1.0, subtitle or 'Generated by Deepernova', 22, style['content_color'], bold=False, align=PP_ALIGN.LEFT)


def generate_ppt(json_data):
    """Generate PowerPoint presentation from JSON"""
    try:
        # Validate input
        is_valid, errors = validate_input(json_data)
        if not is_valid:
            return False, {"error": "Input validation failed", "details": errors}
        
        # Setup output directory
        setup_output_dir()
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        template_name = json_data.get('template', 'classic')
        style = get_template_style(template_name)

        # Add title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_title_slide_style(title_slide, json_data.get('title', 'Untitled'), json_data.get('subtitle', 'Generated by Deepernova'), style)

        # Add content slides
        slides_data = json_data.get('slides', [])
        for slide_data in slides_data:
            slide_title = slide_data.get('title', 'Slide')
            content = slide_data.get('content', '')
            add_slide_with_content(prs, slide_title, content, style)
        
        # Save presentation
        filename = f"presentation_{int(datetime.now().timestamp())}.pptx"
        filepath = OUTPUT_DIR / filename
        prs.save(filepath)
        
        # Check file size
        file_size_mb = filepath.stat().st_size / (1024 * 1024)
        if file_size_mb > MAX_FILE_SIZE_MB:
            filepath.unlink()
            return False, {"error": f"Generated file too large ({file_size_mb:.2f}MB)"}
        
        return True, {
            "success": True,
            "filename": filename,
            "path": str(filepath),
            "size_mb": round(file_size_mb, 2),
            "slides": len(slides_data) + 1
        }
    
    except TimeoutError:
        return False, {"error": "Generation timeout - presentation too complex"}
    except Exception as e:
        return False, {"error": f"Generation failed: {str(e)}"}

def main():
    """Main entry point"""
    try:
        # Read JSON from stdin
        input_data = sys.stdin.read()
        json_data = json.loads(input_data)
        
        # Generate PPT with timeout using threading
        result_container = {'success': False, 'result': None}
        exception_container = {'exception': None}
        
        def generate_with_timeout():
            try:
                success, result = generate_ppt(json_data)
                result_container['success'] = success
                result_container['result'] = result
            except Exception as e:
                exception_container['exception'] = e
        
        # Run generation in thread with timeout
        gen_thread = threading.Thread(target=generate_with_timeout, daemon=False)
        gen_thread.start()
        gen_thread.join(timeout=TIMEOUT_SECONDS)
        
        # Check if thread is still alive (timeout occurred)
        if gen_thread.is_alive():
            # Thread still running - timeout exceeded
            print(json.dumps({
                "success": False,
                "data": {"error": "Generation timeout - presentation too complex"}
            }))
            sys.exit(1)
        
        # Check for exceptions
        if exception_container['exception']:
            raise exception_container['exception']
        
        # Output result as JSON
        print(json.dumps({
            "success": result_container['success'],
            "data": result_container['result']
        }))
        
        sys.exit(0 if result_container['success'] else 1)
    
    except json.JSONDecodeError:
        print(json.dumps({"success": False, "data": {"error": "Invalid JSON input"}}))
        sys.exit(1)
    except Exception as e:
        print(json.dumps({"success": False, "data": {"error": f"Unexpected error: {str(e)}"}}))
        sys.exit(1)

if __name__ == '__main__':
    main()
